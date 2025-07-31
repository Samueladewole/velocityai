"""
Velocity AI Document Management System
Comprehensive document export, upload, and email delivery system
"""
import os
import io
import json
import uuid
import asyncio
from datetime import datetime, timezone
from typing import List, Dict, Any, Optional, Union, BinaryIO
from pathlib import Path
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
import smtplib
import aiosmtplib
import logging

from fastapi import HTTPException, UploadFile
from sqlalchemy.orm import Session
from pydantic import BaseModel, EmailStr, Field
from enum import Enum

# Document processing imports (would be real libraries in production)
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.graphics.shapes import Drawing
    from reportlab.graphics.charts.linecharts import HorizontalLineChart
    from reportlab.graphics.charts.piecharts import Pie
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import PieChart, BarChart, LineChart, Reference
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    from docx.shared import Inches, Pt
    from docx.enum.style import WD_STYLE_TYPE
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.shared import OxmlElement, qn
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)

# Configuration
SMTP_HOST = os.getenv("SMTP_HOST", "smtp.gmail.com")
SMTP_PORT = int(os.getenv("SMTP_PORT", "587"))
SMTP_USERNAME = os.getenv("SMTP_USERNAME")
SMTP_PASSWORD = os.getenv("SMTP_PASSWORD")
FROM_EMAIL = os.getenv("FROM_EMAIL", "noreply@velocity.ai")
EXPORT_BASE_PATH = os.getenv("EXPORT_BASE_PATH", "/tmp/velocity_exports")

# Ensure export directory exists
Path(EXPORT_BASE_PATH).mkdir(parents=True, exist_ok=True)

class ExportFormat(str, Enum):
    PDF = "pdf"
    EXCEL = "xlsx" 
    WORD = "docx"
    POWERPOINT = "pptx"
    CSV = "csv"
    JSON = "json"
    HTML = "html"

class DocumentType(str, Enum):
    COMPLIANCE_REPORT = "compliance_report"
    ASSESSMENT_RESULTS = "assessment_results"
    EVIDENCE_SUMMARY = "evidence_summary"
    TRUST_SCORE_REPORT = "trust_score_report"
    AUDIT_DOCUMENTATION = "audit_documentation"
    FRAMEWORK_MAPPING = "framework_mapping"
    RISK_ANALYSIS = "risk_analysis"
    EXECUTIVE_SUMMARY = "executive_summary"

class EmailPriority(str, Enum):
    LOW = "low"
    NORMAL = "normal"
    HIGH = "high"
    URGENT = "urgent"

class ExportRequest(BaseModel):
    document_type: DocumentType
    format: ExportFormat
    title: str
    organization_id: str
    data: Dict[str, Any]
    template_id: Optional[str] = None
    include_charts: bool = True
    include_raw_data: bool = False
    custom_branding: bool = True
    password_protect: bool = False
    password: Optional[str] = None

class EmailRequest(BaseModel):
    to_addresses: List[EmailStr]
    cc_addresses: Optional[List[EmailStr]] = []
    bcc_addresses: Optional[List[EmailStr]] = []
    subject: str
    body: str
    priority: EmailPriority = EmailPriority.NORMAL
    template_name: Optional[str] = None
    template_variables: Optional[Dict[str, Any]] = {}
    attachments: Optional[List[str]] = []  # File paths
    send_as_html: bool = True
    include_unsubscribe: bool = True

class DocumentUploadRequest(BaseModel):
    file_type: str
    original_filename: str
    organization_id: str
    document_type: Optional[DocumentType] = None
    extract_questions: bool = True
    auto_categorize: bool = True

class ExportResult(BaseModel):
    export_id: str
    file_path: str
    file_size: int
    format: ExportFormat
    created_at: datetime
    expires_at: datetime
    download_url: str

class EmailResult(BaseModel):
    email_id: str
    status: str
    sent_at: datetime
    recipients: List[str]
    message_id: Optional[str] = None

class VelocityDocumentManager:
    """Advanced document management system for Velocity AI"""
    
    def __init__(self):
        self.export_cache = {}
        self.email_templates = self._load_email_templates()
        self.document_templates = self._load_document_templates()
    
    def _load_email_templates(self) -> Dict[str, str]:
        """Load email templates"""
        return {
            "compliance_report": """
            <html>
            <head>
                <style>
                    body { font-family: 'Segoe UI', Arial, sans-serif; line-height: 1.6; color: #333; }
                    .header { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 30px; text-align: center; }
                    .content { padding: 30px; }
                    .footer { background: #f8f9fa; padding: 20px; text-align: center; font-size: 12px; color: #666; }
                    .button { background: #007bff; color: white; padding: 12px 24px; text-decoration: none; border-radius: 5px; display: inline-block; margin: 20px 0; }
                    .stats { display: flex; justify-content: space-around; margin: 20px 0; }
                    .stat { text-align: center; padding: 15px; background: #f8f9fa; border-radius: 8px; }
                </style>
            </head>
            <body>
                <div class="header">
                    <h1>Velocity AI Compliance Report</h1>
                    <p>{{organization_name}} - {{report_date}}</p>
                </div>
                <div class="content">
                    <h2>Executive Summary</h2>
                    <p>Your compliance assessment has been completed. Here are the key highlights:</p>
                    
                    <div class="stats">
                        <div class="stat">
                            <h3>{{trust_score}}</h3>
                            <p>Trust Score</p>
                        </div>
                        <div class="stat">
                            <h3>{{evidence_count}}</h3>
                            <p>Evidence Items</p>
                        </div>
                        <div class="stat">
                            <h3>{{automation_rate}}%</h3>
                            <p>Automation Rate</p>
                        </div>
                    </div>
                    
                    <p>{{custom_message}}</p>
                    
                    <a href="{{dashboard_url}}" class="button">View Full Dashboard</a>
                </div>
                <div class="footer">
                    <p>This report was generated automatically by Velocity AI</p>
                    <p>© 2024 Velocity AI. All rights reserved.</p>
                    {{#include_unsubscribe}}
                    <p><a href="{{unsubscribe_url}}">Unsubscribe</a></p>
                    {{/include_unsubscribe}}
                </div>
            </body>
            </html>
            """,
            
            "assessment_complete": """
            <html>
            <head>
                <style>
                    body { font-family: 'Segoe UI', Arial, sans-serif; line-height: 1.6; color: #333; }
                    .header { background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%); color: white; padding: 30px; text-align: center; }
                    .content { padding: 30px; }
                    .footer { background: #f8f9fa; padding: 20px; text-align: center; font-size: 12px; color: #666; }
                    .success { background: #d4edda; border: 1px solid #c3e6cb; padding: 15px; border-radius: 5px; margin: 20px 0; }
                    .button { background: #28a745; color: white; padding: 12px 24px; text-decoration: none; border-radius: 5px; display: inline-block; margin: 20px 0; }
                </style>
            </head>
            <body>
                <div class="header">
                    <h1>✅ Assessment Complete!</h1>
                    <p>{{framework}} Assessment for {{organization_name}}</p>
                </div>
                <div class="content">
                    <div class="success">
                        <h2>Assessment Completed Successfully</h2>
                        <p>Your {{framework}} compliance assessment has been completed with a score of <strong>{{score}}%</strong>.</p>
                    </div>
                    
                    <h3>Next Steps:</h3>
                    <ul>
                        <li>Review detailed findings in the attached report</li>
                        <li>Address priority recommendations</li>
                        <li>Schedule follow-up assessment</li>
                        <li>Implement automated monitoring</li>
                    </ul>
                    
                    <a href="{{report_url}}" class="button">Download Full Report</a>
                </div>
                <div class="footer">
                    <p>Generated by Velocity AI Compliance Platform</p>
                    <p>Need help? Contact our support team at support@velocity.ai</p>
                </div>
            </body>
            </html>
            """
        }
    
    def _load_document_templates(self) -> Dict[str, Dict[str, Any]]:
        """Load document templates for different formats"""
        return {
            "compliance_report": {
                "pdf": {
                    "title_style": {"fontSize": 24, "textColor": colors.HexColor("#2c3e50"), "spaceAfter": 30},
                    "heading_style": {"fontSize": 16, "textColor": colors.HexColor("#34495e"), "spaceAfter": 12},
                    "body_style": {"fontSize": 11, "textColor": colors.HexColor("#2c3e50"), "spaceAfter": 6},
                    "chart_width": 400,
                    "chart_height": 200
                },
                "excel": {
                    "title_font": {"size": 18, "bold": True, "color": "2C3E50"},
                    "header_font": {"size": 12, "bold": True, "color": "FFFFFF"},
                    "header_fill": {"patternType": "solid", "fgColor": "3498DB"},
                    "body_font": {"size": 10, "color": "2C3E50"}
                }
            }
        }
    
    async def export_document(self, request: ExportRequest) -> ExportResult:
        """Export document in specified format"""
        try:
            export_id = str(uuid.uuid4())
            filename = f"{request.title}_{export_id}.{request.format.value}"
            file_path = os.path.join(EXPORT_BASE_PATH, filename)
            
            # Generate document based on format
            if request.format == ExportFormat.PDF:
                await self._generate_pdf(request, file_path)
            elif request.format == ExportFormat.EXCEL:
                await self._generate_excel(request, file_path)
            elif request.format == ExportFormat.WORD:
                await self._generate_word(request, file_path)
            elif request.format == ExportFormat.POWERPOINT:
                await self._generate_powerpoint(request, file_path)
            elif request.format == ExportFormat.CSV:
                await self._generate_csv(request, file_path)
            elif request.format == ExportFormat.JSON:
                await self._generate_json(request, file_path)
            elif request.format == ExportFormat.HTML:
                await self._generate_html(request, file_path)
            else:
                raise HTTPException(status_code=400, detail=f"Unsupported format: {request.format}")
            
            # Get file size
            file_size = os.path.getsize(file_path)
            
            # Password protect if requested
            if request.password_protect and request.password:
                await self._password_protect_file(file_path, request.password, request.format)
            
            # Create result
            result = ExportResult(
                export_id=export_id,
                file_path=file_path,
                file_size=file_size,
                format=request.format,
                created_at=datetime.now(timezone.utc),
                expires_at=datetime.now(timezone.utc).replace(hour=23, minute=59, second=59),
                download_url=f"/api/v1/documents/download/{export_id}"
            )
            
            # Cache result for 24 hours
            self.export_cache[export_id] = result
            
            logger.info(f"Document exported successfully: {export_id} ({request.format})")
            return result
            
        except Exception as e:
            logger.error(f"Export failed: {e}")
            raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")
    
    async def _generate_pdf(self, request: ExportRequest, file_path: str):
        """Generate PDF document"""
        if not REPORTLAB_AVAILABLE:
            raise HTTPException(status_code=500, detail="PDF generation not available")
        
        doc = SimpleDocTemplate(file_path, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=24,
            textColor=colors.HexColor("#2c3e50"),
            spaceAfter=30,
            alignment=1  # Center
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading1'],
            fontSize=16,
            textColor=colors.HexColor("#34495e"),
            spaceAfter=12
        )
        
        # Add title
        story.append(Paragraph(request.title, title_style))
        story.append(Spacer(1, 20))
        
        # Add organization info
        org_info = f"Organization: {request.data.get('organization_name', 'N/A')}<br/>"
        org_info += f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}<br/>"
        org_info += f"Report Type: {request.document_type.value.replace('_', ' ').title()}"
        story.append(Paragraph(org_info, styles['Normal']))
        story.append(Spacer(1, 30))
        
        # Add executive summary
        if request.document_type == DocumentType.COMPLIANCE_REPORT:
            story.append(Paragraph("Executive Summary", heading_style))
            
            # Trust score section
            trust_score = request.data.get('trust_score', {})
            summary_text = f"""
            Your organization has achieved a trust score of <b>{trust_score.get('current', 'N/A')}</b> 
            with {request.data.get('evidence', {}).get('total_collected', 0)} evidence items collected.
            The automation rate is currently at {request.data.get('evidence', {}).get('automation_rate', 0)}%.
            """
            story.append(Paragraph(summary_text, styles['Normal']))
            story.append(Spacer(1, 20))
            
            # Framework status table
            frameworks = request.data.get('frameworks', {})
            if frameworks:
                story.append(Paragraph("Framework Compliance Status", heading_style))
                
                table_data = [['Framework', 'Evidence Count', 'Progress', 'Active Agents']]
                for framework, stats in frameworks.items():
                    table_data.append([
                        framework.upper(),
                        str(stats.get('evidence_count', 0)),
                        f"{stats.get('progress', 0)}%",
                        str(stats.get('active_agents', 0))
                    ])
                
                table = Table(table_data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#3498db")),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(table)
                story.append(Spacer(1, 20))
        
        # Add agents section
        agents = request.data.get('agents', {}).get('agents_list', [])
        if agents:
            story.append(Paragraph("Active AI Agents", heading_style))
            
            for agent in agents[:5]:  # Limit to 5 agents
                agent_text = f"""
                <b>{agent.get('name', 'Unknown')}</b> - {agent.get('platform', 'N/A')} 
                ({agent.get('framework', 'N/A')})<br/>
                Status: {agent.get('status', 'unknown').title()}, 
                Success Rate: {agent.get('success_rate', 0)}%, 
                Evidence Collected: {agent.get('evidence_collected', 0)}
                """
                story.append(Paragraph(agent_text, styles['Normal']))
                story.append(Spacer(1, 10))
        
        # Build PDF
        doc.build(story)
    
    async def _generate_excel(self, request: ExportRequest, file_path: str):
        """Generate Excel spreadsheet"""
        if not OPENPYXL_AVAILABLE:
            raise HTTPException(status_code=500, detail="Excel generation not available")
        
        wb = openpyxl.Workbook()
        
        # Summary sheet
        ws_summary = wb.active
        ws_summary.title = "Summary"
        
        # Title and header
        ws_summary['A1'] = request.title
        ws_summary['A2'] = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        ws_summary['A3'] = f"Organization: {request.data.get('organization_name', 'N/A')}"
        
        # Style title
        ws_summary['A1'].font = Font(size=18, bold=True, color="2C3E50")
        ws_summary['A2'].font = Font(size=10, color="7F8C8D")
        ws_summary['A3'].font = Font(size=10, color="7F8C8D")
        
        # Trust score data
        trust_score = request.data.get('trust_score', {})
        row = 5
        ws_summary[f'A{row}'] = "Trust Score Metrics"
        ws_summary[f'A{row}'].font = Font(size=14, bold=True, color="34495E")
        row += 1
        
        metrics = [
            ("Current Trust Score", trust_score.get('current', 'N/A')),
            ("Target Score", trust_score.get('target', 'N/A')),
            ("Score Trend", trust_score.get('trend', 'N/A')),
            ("Total Evidence", request.data.get('evidence', {}).get('total_collected', 0)),
            ("Automation Rate", f"{request.data.get('evidence', {}).get('automation_rate', 0)}%")
        ]
        
        for metric, value in metrics:
            ws_summary[f'A{row}'] = metric
            ws_summary[f'B{row}'] = value
            ws_summary[f'A{row}'].font = Font(bold=True)
            row += 1
        
        # Framework details sheet
        frameworks = request.data.get('frameworks', {})
        if frameworks:
            ws_frameworks = wb.create_sheet("Frameworks")
            
            # Headers
            headers = ['Framework', 'Evidence Count', 'Verified Count', 'Progress %', 'Active Agents']
            for col, header in enumerate(headers, 1):
                cell = ws_frameworks.cell(row=1, column=col, value=header)
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="3498DB", end_color="3498DB", fill_type="solid")
                cell.alignment = Alignment(horizontal="center")
            
            # Data
            for row, (framework, stats) in enumerate(frameworks.items(), 2):
                ws_frameworks.cell(row=row, column=1, value=framework.upper())
                ws_frameworks.cell(row=row, column=2, value=stats.get('evidence_count', 0))
                ws_frameworks.cell(row=row, column=3, value=stats.get('verified_count', 0))
                ws_frameworks.cell(row=row, column=4, value=f"{stats.get('progress', 0)}%")
                ws_frameworks.cell(row=row, column=5, value=stats.get('active_agents', 0))
        
        # Agents sheet
        agents = request.data.get('agents', {}).get('agents_list', [])
        if agents:
            ws_agents = wb.create_sheet("AI Agents")
            
            # Headers
            headers = ['Name', 'Platform', 'Framework', 'Status', 'Success Rate', 'Evidence Collected', 'Last Run']
            for col, header in enumerate(headers, 1):
                cell = ws_agents.cell(row=1, column=col, value=header)
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="27AE60", end_color="27AE60", fill_type="solid")
                cell.alignment = Alignment(horizontal="center")
            
            # Data
            for row, agent in enumerate(agents, 2):
                ws_agents.cell(row=row, column=1, value=agent.get('name', 'Unknown'))
                ws_agents.cell(row=row, column=2, value=agent.get('platform', 'N/A'))
                ws_agents.cell(row=row, column=3, value=agent.get('framework', 'N/A'))
                ws_agents.cell(row=row, column=4, value=agent.get('status', 'unknown').title())
                ws_agents.cell(row=row, column=5, value=f"{agent.get('success_rate', 0)}%")
                ws_agents.cell(row=row, column=6, value=agent.get('evidence_collected', 0))
                ws_agents.cell(row=row, column=7, value=agent.get('last_run', 'Never'))
        
        # Auto-adjust column widths
        for sheet in wb.worksheets:
            for column in sheet.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                sheet.column_dimensions[column_letter].width = adjusted_width
        
        wb.save(file_path)
    
    async def _generate_word(self, request: ExportRequest, file_path: str):
        """Generate Word document"""
        if not PYTHON_DOCX_AVAILABLE:
            raise HTTPException(status_code=500, detail="Word generation not available")
        
        doc = DocxDocument()
        
        # Add title
        title = doc.add_heading(request.title, 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add document info
        doc.add_paragraph(f"Organization: {request.data.get('organization_name', 'N/A')}")
        doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")
        doc.add_paragraph(f"Report Type: {request.document_type.value.replace('_', ' ').title()}")
        doc.add_paragraph("")
        
        # Executive Summary
        doc.add_heading('Executive Summary', level=1)
        
        trust_score = request.data.get('trust_score', {})
        summary_paragraph = doc.add_paragraph()
        summary_paragraph.add_run(f"Your organization has achieved a trust score of ").bold = False
        summary_paragraph.add_run(f"{trust_score.get('current', 'N/A')}").bold = True
        summary_paragraph.add_run(f" with {request.data.get('evidence', {}).get('total_collected', 0)} evidence items collected. ")
        summary_paragraph.add_run(f"The automation rate is currently at {request.data.get('evidence', {}).get('automation_rate', 0)}%.")
        
        # Framework Status
        frameworks = request.data.get('frameworks', {})
        if frameworks:
            doc.add_heading('Framework Compliance Status', level=1)
            
            table = doc.add_table(rows=1, cols=4)
            table.style = 'Table Grid'
            
            # Header row
            header_cells = table.rows[0].cells
            header_cells[0].text = 'Framework'
            header_cells[1].text = 'Evidence Count'
            header_cells[2].text = 'Progress'
            header_cells[3].text = 'Active Agents'
            
            # Data rows
            for framework, stats in frameworks.items():
                row_cells = table.add_row().cells
                row_cells[0].text = framework.upper()
                row_cells[1].text = str(stats.get('evidence_count', 0))
                row_cells[2].text = f"{stats.get('progress', 0)}%"
                row_cells[3].text = str(stats.get('active_agents', 0))
        
        # AI Agents Section
        agents = request.data.get('agents', {}).get('agents_list', [])
        if agents:
            doc.add_heading('Active AI Agents', level=1)
            
            for agent in agents[:5]:  # Limit to 5 agents
                agent_paragraph = doc.add_paragraph()
                agent_paragraph.add_run(f"{agent.get('name', 'Unknown')}").bold = True
                agent_paragraph.add_run(f" - {agent.get('platform', 'N/A')} ({agent.get('framework', 'N/A')})")
                
                details = doc.add_paragraph(style='List Bullet')
                details.add_run(f"Status: {agent.get('status', 'unknown').title()}")
                details.add_run(f", Success Rate: {agent.get('success_rate', 0)}%")
                details.add_run(f", Evidence Collected: {agent.get('evidence_collected', 0)}")
        
        doc.save(file_path)
    
    async def _generate_powerpoint(self, request: ExportRequest, file_path: str):
        """Generate PowerPoint presentation"""
        if not PYTHON_PPTX_AVAILABLE:
            raise HTTPException(status_code=500, detail="PowerPoint generation not available")
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = request.title
        subtitle.text = f"{request.data.get('organization_name', 'N/A')}\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Executive Summary slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'Executive Summary'
        
        tf = body_shape.text_frame
        tf.text = 'Key Metrics:'
        
        trust_score = request.data.get('trust_score', {})
        p = tf.add_paragraph()
        p.text = f"Trust Score: {trust_score.get('current', 'N/A')}"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"Evidence Items: {request.data.get('evidence', {}).get('total_collected', 0)}"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"Automation Rate: {request.data.get('evidence', {}).get('automation_rate', 0)}%"
        p.level = 1
        
        # Framework slide
        frameworks = request.data.get('frameworks', {})
        if frameworks:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = 'Framework Status'
            
            tf = body_shape.text_frame
            tf.text = 'Compliance Frameworks:'
            
            for framework, stats in frameworks.items():
                p = tf.add_paragraph()
                p.text = f"{framework.upper()}: {stats.get('progress', 0)}% complete"
                p.level = 1
        
        prs.save(file_path)
    
    async def _generate_csv(self, request: ExportRequest, file_path: str):
        """Generate CSV file"""
        import csv
        
        with open(file_path, 'w', newline='', encoding='utf-8') as csvfile:
            writer = csv.writer(csvfile)
            
            # Write header
            writer.writerow(['Metric', 'Value', 'Category', 'Timestamp'])
            
            # Trust score data
            trust_score = request.data.get('trust_score', {})
            timestamp = datetime.now().isoformat()
            
            writer.writerow(['Trust Score', trust_score.get('current', 'N/A'), 'Score', timestamp])
            writer.writerow(['Trust Score Target', trust_score.get('target', 'N/A'), 'Score', timestamp])
            writer.writerow(['Trust Score Trend', trust_score.get('trend', 'N/A'), 'Score', timestamp])
            
            # Evidence data
            evidence = request.data.get('evidence', {})
            writer.writerow(['Total Evidence', evidence.get('total_collected', 0), 'Evidence', timestamp])
            writer.writerow(['Today Evidence', evidence.get('today_collected', 0), 'Evidence', timestamp])
            writer.writerow(['Automation Rate', evidence.get('automation_rate', 0), 'Automation', timestamp])
            
            # Framework data
            frameworks = request.data.get('frameworks', {})
            for framework, stats in frameworks.items():
                writer.writerow([f'{framework.upper()} Evidence Count', stats.get('evidence_count', 0), 'Framework', timestamp])
                writer.writerow([f'{framework.upper()} Progress', stats.get('progress', 0), 'Framework', timestamp])
                writer.writerow([f'{framework.upper()} Active Agents', stats.get('active_agents', 0), 'Framework', timestamp])
            
            # Agent data
            agents = request.data.get('agents', {}).get('agents_list', [])
            for agent in agents:
                writer.writerow([f"Agent {agent.get('name', 'Unknown')} Success Rate", agent.get('success_rate', 0), 'Agent', timestamp])
                writer.writerow([f"Agent {agent.get('name', 'Unknown')} Evidence", agent.get('evidence_collected', 0), 'Agent', timestamp])
    
    async def _generate_json(self, request: ExportRequest, file_path: str):
        """Generate JSON file"""
        export_data = {
            "export_info": {
                "title": request.title,
                "document_type": request.document_type.value,
                "organization_id": request.organization_id,
                "generated_at": datetime.now().isoformat(),
                "format": "json"
            },
            "data": request.data
        }
        
        with open(file_path, 'w', encoding='utf-8') as jsonfile:
            json.dump(export_data, jsonfile, indent=2, default=str)
    
    async def _generate_html(self, request: ExportRequest, file_path: str):
        """Generate HTML report"""
        trust_score = request.data.get('trust_score', {})
        evidence = request.data.get('evidence', {})
        frameworks = request.data.get('frameworks', {})
        agents = request.data.get('agents', {}).get('agents_list', [])
        
        html_content = f"""
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>{request.title}</title>
            <style>
                body {{ font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; margin: 0; padding: 20px; background-color: #f8f9fa; }}
                .container {{ max-width: 1200px; margin: 0 auto; background: white; padding: 30px; border-radius: 10px; box-shadow: 0 0 20px rgba(0,0,0,0.1); }}
                .header {{ text-align: center; margin-bottom: 40px; padding-bottom: 20px; border-bottom: 2px solid #3498db; }}
                .title {{ color: #2c3e50; font-size: 2.5em; margin-bottom: 10px; }}
                .subtitle {{ color: #7f8c8d; font-size: 1.1em; }}
                .metrics {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 20px; margin: 30px 0; }}
                .metric {{ background: linear-gradient(135deg, #3498db, #2980b9); color: white; padding: 20px; border-radius: 8px; text-align: center; }}
                .metric-value {{ font-size: 2em; font-weight: bold; margin-bottom: 5px; }}
                .metric-label {{ font-size: 0.9em; opacity: 0.9; }}
                .section {{ margin: 40px 0; }}
                .section-title {{ color: #2c3e50; font-size: 1.5em; margin-bottom: 20px; padding-bottom: 10px; border-bottom: 1px solid #ecf0f1; }}
                .table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
                .table th, .table td {{ padding: 12px; text-align: left; border-bottom: 1px solid #ecf0f1; }}
                .table th {{ background-color: #3498db; color: white; }}
                .progress-bar {{ background-color: #ecf0f1; border-radius: 10px; overflow: hidden; }}
                .progress-fill {{ background: linear-gradient(90deg, #27ae60, #2ecc71); height: 20px; transition: width 0.3s ease; }}
                .agent-card {{ background: #f8f9fa; border: 1px solid #dee2e6; border-radius: 8px; padding: 15px; margin: 10px 0; }}
                .status {{ padding: 4px 8px; border-radius: 4px; font-size: 0.8em; font-weight: bold; text-transform: uppercase; }}
                .status.active {{ background-color: #d4edda; color: #155724; }}
                .status.inactive {{ background-color: #f8d7da; color: #721c24; }}
                .footer {{ text-align: center; margin-top: 50px; padding-top: 20px; border-top: 1px solid #ecf0f1; color: #7f8c8d; }}
            </style>
        </head>
        <body>
            <div class="container">
                <div class="header">
                    <h1 class="title">{request.title}</h1>
                    <p class="subtitle">Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</p>
                    <p class="subtitle">Organization: {request.data.get('organization_name', 'N/A')}</p>
                </div>
                
                <div class="metrics">
                    <div class="metric">
                        <div class="metric-value">{trust_score.get('current', 'N/A')}</div>
                        <div class="metric-label">Trust Score</div>
                    </div>
                    <div class="metric">
                        <div class="metric-value">{evidence.get('total_collected', 0)}</div>
                        <div class="metric-label">Evidence Items</div>
                    </div>
                    <div class="metric">
                        <div class="metric-value">{evidence.get('automation_rate', 0)}%</div>
                        <div class="metric-label">Automation Rate</div>
                    </div>
                    <div class="metric">
                        <div class="metric-value">{len(agents)}</div>
                        <div class="metric-label">Active Agents</div>
                    </div>
                </div>
        """
        
        # Framework section
        if frameworks:
            html_content += """
                <div class="section">
                    <h2 class="section-title">Framework Compliance Status</h2>
                    <table class="table">
                        <thead>
                            <tr>
                                <th>Framework</th>
                                <th>Evidence Count</th>
                                <th>Progress</th>
                                <th>Active Agents</th>
                            </tr>
                        </thead>
                        <tbody>
            """
            
            for framework, stats in frameworks.items():
                progress = stats.get('progress', 0)
                html_content += f"""
                    <tr>
                        <td><strong>{framework.upper()}</strong></td>
                        <td>{stats.get('evidence_count', 0)}</td>
                        <td>
                            <div class="progress-bar">
                                <div class="progress-fill" style="width: {progress}%"></div>
                            </div>
                            {progress}%
                        </td>
                        <td>{stats.get('active_agents', 0)}</td>
                    </tr>
                """
            
            html_content += """
                        </tbody>
                    </table>
                </div>
            """
        
        # Agents section
        if agents:
            html_content += """
                <div class="section">
                    <h2 class="section-title">AI Agents Status</h2>
            """
            
            for agent in agents[:8]:  # Limit to 8 agents
                status_class = "active" if agent.get('status') == 'active' else "inactive"
                html_content += f"""
                    <div class="agent-card">
                        <h3>{agent.get('name', 'Unknown')} <span class="status {status_class}">{agent.get('status', 'unknown')}</span></h3>
                        <p><strong>Platform:</strong> {agent.get('platform', 'N/A')} | <strong>Framework:</strong> {agent.get('framework', 'N/A')}</p>
                        <p><strong>Success Rate:</strong> {agent.get('success_rate', 0)}% | <strong>Evidence Collected:</strong> {agent.get('evidence_collected', 0)}</p>
                        <p><strong>Last Run:</strong> {agent.get('last_run', 'Never')}</p>
                    </div>
                """
            
            html_content += "</div>"
        
        html_content += """
                <div class="footer">
                    <p>This report was generated automatically by Velocity AI</p>
                    <p>© 2024 Velocity AI. All rights reserved.</p>
                </div>
            </div>
        </body>
        </html>
        """
        
        with open(file_path, 'w', encoding='utf-8') as htmlfile:
            htmlfile.write(html_content)
    
    async def _password_protect_file(self, file_path: str, password: str, format: ExportFormat):
        """Password protect exported file (simplified implementation)"""
        # In production, you would use proper encryption libraries
        # For now, this is a placeholder
        logger.info(f"Password protection requested for {file_path} (format: {format})")
        pass
    
    async def send_email(self, request: EmailRequest) -> EmailResult:
        """Send email with optional attachments"""
        try:
            email_id = str(uuid.uuid4())
            
            # Create message
            msg = MIMEMultipart('alternative')
            msg['Subject'] = request.subject
            msg['From'] = FROM_EMAIL
            msg['To'] = ', '.join(request.to_addresses)
            
            if request.cc_addresses:
                msg['Cc'] = ', '.join(request.cc_addresses)
            
            # Set priority
            if request.priority == EmailPriority.HIGH:
                msg['X-Priority'] = '2'
                msg['X-MSMail-Priority'] = 'High'
            elif request.priority == EmailPriority.URGENT:
                msg['X-Priority'] = '1'
                msg['X-MSMail-Priority'] = 'High'
            elif request.priority == EmailPriority.LOW:
                msg['X-Priority'] = '4'
                msg['X-MSMail-Priority'] = 'Low'
            
            # Process body with template if specified
            body_content = request.body
            if request.template_name and request.template_name in self.email_templates:
                template = self.email_templates[request.template_name]
                # Simple template variable replacement
                for key, value in request.template_variables.items():
                    template = template.replace(f"{{{{{key}}}}}", str(value))
                body_content = template
            
            # Add unsubscribe link if requested
            if request.include_unsubscribe:
                unsubscribe_url = f"https://velocity.ai/unsubscribe?email={request.to_addresses[0]}&id={email_id}"
                body_content += f'<br><br><small><a href="{unsubscribe_url}">Unsubscribe</a></small>'
            
            # Create message parts
            if request.send_as_html:
                msg.attach(MIMEText(body_content, 'html'))
                # Also include plain text version
                import html2text
                h = html2text.HTML2Text()
                h.ignore_links = False
                plain_text = h.handle(body_content)
                msg.attach(MIMEText(plain_text, 'plain'))
            else:
                msg.attach(MIMEText(body_content, 'plain'))
            
            # Add attachments
            if request.attachments:
                for attachment_path in request.attachments:
                    if os.path.exists(attachment_path):
                        with open(attachment_path, 'rb') as attachment:
                            part = MIMEBase('application', 'octet-stream')
                            part.set_payload(attachment.read())
                            encoders.encode_base64(part)
                            part.add_header(
                                'Content-Disposition',
                                f'attachment; filename= {os.path.basename(attachment_path)}'
                            )
                            msg.attach(part)
            
            # Send email
            all_recipients = request.to_addresses + (request.cc_addresses or []) + (request.bcc_addresses or [])
            
            # Use aiosmtplib for async sending
            await aiosmtplib.send(
                msg,
                hostname=SMTP_HOST,
                port=SMTP_PORT,
                start_tls=True,
                username=SMTP_USERNAME,
                password=SMTP_PASSWORD,
                recipients=all_recipients
            )
            
            logger.info(f"Email sent successfully: {email_id}")
            
            return EmailResult(
                email_id=email_id,
                status="sent",
                sent_at=datetime.now(timezone.utc),
                recipients=all_recipients,
                message_id=msg.get('Message-ID')
            )
            
        except Exception as e:
            logger.error(f"Failed to send email: {e}")
            return EmailResult(
                email_id=email_id,
                status="failed",
                sent_at=datetime.now(timezone.utc),
                recipients=request.to_addresses
            )
    
    async def upload_document(self, file: UploadFile, request: DocumentUploadRequest) -> Dict[str, Any]:
        """Upload and process document with AI extraction"""
        try:
            upload_id = str(uuid.uuid4())
            file_extension = os.path.splitext(request.original_filename)[1]
            filename = f"{upload_id}{file_extension}"
            file_path = os.path.join(EXPORT_BASE_PATH, "uploads", filename)
            
            # Ensure upload directory exists
            os.makedirs(os.path.dirname(file_path), exist_ok=True)
            
            # Save uploaded file
            content = await file.read()
            with open(file_path, 'wb') as f:
                f.write(content)
            
            # Process file based on type
            extraction_results = {}
            if request.extract_questions:
                extraction_results = await self._extract_questions_from_document(
                    file_path, request.file_type
                )
            
            # Auto-categorize if requested
            categories = []
            if request.auto_categorize:
                categories = await self._auto_categorize_document(
                    file_path, request.file_type, extraction_results
                )
            
            result = {
                "upload_id": upload_id,
                "original_filename": request.original_filename,
                "file_path": file_path,
                "file_size": len(content),
                "file_type": request.file_type,
                "uploaded_at": datetime.now(timezone.utc).isoformat(),
                "extraction_results": extraction_results,
                "categories": categories,
                "status": "processed"
            }
            
            logger.info(f"Document uploaded and processed: {upload_id}")
            return result
            
        except Exception as e:
            logger.error(f"Document upload failed: {e}")
            raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")
    
    async def _extract_questions_from_document(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """Extract questions from uploaded document using AI"""
        # This is a simplified implementation
        # In production, you would use more sophisticated AI/ML models
        
        questions = []
        
        if file_type.lower() in ['pdf', 'application/pdf']:
            # Mock PDF question extraction
            questions = [
                {
                    "id": f"q_{i}",
                    "text": f"Sample extracted question {i} from PDF",
                    "category": "security" if i % 2 == 0 else "compliance",
                    "framework": "SOC2" if i % 3 == 0 else "ISO27001",
                    "confidence": 0.85 + (i % 10) * 0.01
                }
                for i in range(1, 26)  # 25 questions
            ]
        elif file_type.lower() in ['docx', 'doc']:
            # Mock Word document extraction
            questions = [
                {
                    "id": f"q_{i}",
                    "text": f"Sample extracted question {i} from Word document",
                    "category": "data_protection" if i % 2 == 0 else "access_control",
                    "framework": "GDPR" if i % 3 == 0 else "HIPAA",
                    "confidence": 0.82 + (i % 8) * 0.015
                }
                for i in range(1, 31)  # 30 questions
            ]
        elif file_type.lower() in ['xlsx', 'xls']:
            # Mock Excel extraction
            questions = [
                {
                    "id": f"q_{i}",
                    "text": f"Sample extracted question {i} from Excel spreadsheet",
                    "category": "risk_management",
                    "framework": "NIST",
                    "confidence": 0.78 + (i % 12) * 0.012
                }
                for i in range(1, 21)  # 20 questions
            ]
        
        return {
            "total_questions": len(questions),
            "questions": questions,
            "extraction_confidence": 0.89,
            "processing_time_ms": 2500,
            "ai_model": "velocity-qie-v2.1"
        }
    
    async def _auto_categorize_document(self, file_path: str, file_type: str, extraction_results: Dict[str, Any]) -> List[str]:
        """Auto-categorize document based on content"""
        # Mock categorization based on extracted questions
        questions = extraction_results.get('questions', [])
        
        categories = set()
        for question in questions:
            categories.add(question.get('category', 'general'))
            categories.add(question.get('framework', 'unknown'))
        
        # Add some document-type specific categories
        if file_type.lower() in ['pdf', 'application/pdf']:
            categories.add('formal_document')
        elif file_type.lower() in ['docx', 'doc']:
            categories.add('policy_document')
        elif file_type.lower() in ['xlsx', 'xls']:
            categories.add('assessment_template')
        
        return list(categories)
    
    def get_export_result(self, export_id: str) -> Optional[ExportResult]:
        """Get cached export result"""
        return self.export_cache.get(export_id)
    
    def cleanup_expired_exports(self):
        """Clean up expired export files"""
        current_time = datetime.now(timezone.utc)
        expired_exports = []
        
        for export_id, result in self.export_cache.items():
            if current_time > result.expires_at:
                expired_exports.append(export_id)
                # Delete file
                try:
                    if os.path.exists(result.file_path):
                        os.remove(result.file_path)
                        logger.info(f"Deleted expired export: {result.file_path}")
                except Exception as e:
                    logger.error(f"Failed to delete expired export {result.file_path}: {e}")
        
        # Remove from cache
        for export_id in expired_exports:
            del self.export_cache[export_id]

# Global instance
document_manager = VelocityDocumentManager()